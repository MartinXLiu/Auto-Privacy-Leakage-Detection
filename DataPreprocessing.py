# !/usr/bin/env python
# -*- coding: utf-8 -*-
# @Time    : 2023-08-03 17:26
# @Author  : Jiaxin Liu
# @File    : DataPreprocessing.py
import os
import re

import time
import docx
import openpyxl
import docx2txt
import zipfile
import pytesseract
import email
import base64

from email.header import decode_header
from copy import deepcopy
from PIL import Image
from pptx import Presentation
import win32com.client as wc
from regipy import RegistryHive
from bs4 import BeautifulSoup


def read_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content


# 文本清洗和预处理
def clean_text(text):
    # （可增量更新）需要根据具体情况编写一些正则表达式或规则来清洗文本
    # 这里只是一个简单示例，去除所有非字母和数字的字符
    # cleaned_text = re.sub(r'[^a-zA-Z0-9]', ' ', text)
    # # 编写正则表达式模式
    #     ip_pattern = r'\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}'
    #     username_pattern = r'(?i)username|user|name|login'
    #     password_pattern = r'(?i)password|pass|pwd'
    #     sensitive_info_pattern = fr'({ip_pattern}|{username_pattern}|{password_pattern})'
    cleaned_text = text
    return cleaned_text


# .hiv文件解析函数 .hiv 文件是 Windows 注册表中的一个二进制文件，用于存储特定注册表项及其关联的数据。
def parse_and_save_hiv(hiv_file_path):
    hive = RegistryHive(hiv_file_path)
    output_list = []

    # 遍历注册表键并添加到输出列表中
    for key in hive.recurse_subkeys():
        output_list.append("Key Path: " + key.path)
        output_list.append("Number of Values: " + str(len(key.values)))

        # 遍历注册表值并添加到输出列表中
        for value in key.values:
            value_name = value.name
            value_type = str(value.value_type)
            value_data = value.value
            output_list.append("Value Name: " + value_name)
            output_list.append("Value Type: " + value_type)
            output_list.append("Value Data: " + str(value_data))
            output_list.append("---")

    return output_list


# 提取图片信息
def extract_all_text_from_image(image_path):
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image, lang='eng+chi_sim')  # 识别英文和中文文本
        return text
    except Exception as e:
        print("Error while extracting text from image:", e)
        return ""


# 解码.eml附件
def decode_mime_word(s):
    parts = decode_header(s)
    decoded_parts = []
    for part, encoding in parts:
        if isinstance(part, bytes):
            decoded_part = part.decode(encoding or 'utf-8')
            decoded_parts.append(decoded_part)
        else:
            decoded_parts.append(part)
    return ''.join(decoded_parts)


def preprocess_files_and_save(directory, output_file, preprocess_log):
    with open(output_file, 'w', encoding='utf-8') as output, \
            open(preprocess_log, 'w', encoding='utf-8') as log:
        for root, _, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)

                # 先解压缩文件夹下压缩包，再遍历文件，遍历一次即可
                if file_path.endswith('.zip'):
                    with zipfile.ZipFile(file_path, 'r') as zip_ref:
                        unzip_dir = os.path.join(root, 'unzipped_' + file)
                        zip_ref.extractall(unzip_dir)
                        log.write("Unzipped: " + file_path + " to " + unzip_dir + '\n')

                # 对于.txt文件，直接读取文本内容并进行清洗
                if file_path.endswith('.txt'):
                    with open(file_path, 'r', encoding='utf-8') as f:  # 读取文件
                        text = f.read()
                    cleaned_text = clean_text(text)
                    output.write(cleaned_text + '\n')  # 写入preprocessed_data.txt
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                # 对于.docx文档，使用docx2txt库进行解析并清洗，注意这里需要绝对路径，相对路径会报错
                elif file_path.endswith('.docx'):
                    file_path = os.path.abspath(file_path)  # 绝对路径
                    text = docx2txt.process(file_path)
                    cleaned_text = clean_text(text)
                    output.write(cleaned_text + '\n')
                    with open(file_path.replace('.docx', '.txt'), 'w', encoding='utf-8') as f:
                        f.write(cleaned_text)
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                # 对于.doc文件，使用pywin32库读取文本内容并进行清洗，注意这里需要绝对路径，相对路径会报错
                elif file_path.endswith('.doc'):
                    file_path = os.path.abspath(file_path)  # 绝对路径
                    word = wc.Dispatch('Word.Application')
                    doc = word.Documents.Open(file_path)
                    text = doc.Content.Text
                    doc.Close()
                    word.Quit()
                    cleaned_text = clean_text(text)
                    output.write(cleaned_text + '\n')
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                # 对于.ppt文件，使用pywin32库读取幻灯片内容并进行清洗，注意这里需要绝对路径，相对路径会报错
                # .dps文件，是wps里的ppt
                elif file_path.endswith('.ppt') or file_path.endswith('.dps'):
                    file_path = os.path.abspath(file_path)  # 绝对路径
                    powerpoint = wc.Dispatch("PowerPoint.Application")
                    presentation = powerpoint.Presentations.Open(file_path)
                    text = ""
                    for slide in presentation.Slides:
                        for shape in slide.Shapes:
                            if hasattr(shape, 'HasTextFrame') and shape.HasTextFrame: # 先检查 shape 对象是否具有 HasTextFrame 属性，再检查该属性的值是否为 True
                                text += shape.TextFrame.TextRange.Text + "\n"
                    presentation.Close()  # 只是关闭ppt里面的内容，但是不关闭ppt本身
                    cleaned_text = clean_text(text)
                    output.write(cleaned_text + '\n')
                    log.write("\nPreprocess " + file + " done \n") # 打印日志

                # 对于.pptx文件，使用python-pptx库读取幻灯片内容并进行清洗，注意这里需要绝对路径，相对路径会报错
                elif file_path.endswith('.pptx'):
                    file_path = os.path.abspath(file_path)  # 绝对路径
                    prs = Presentation(file_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text += shape.text + "\n"
                    cleaned_text = clean_text(text)
                    output.write(cleaned_text + '\n')
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                # # 处理.xlsx文件和.et文件
                elif file_path.endswith('.xlsx') or file_path.endswith('.et'):
                    file_path = os.path.abspath(file_path)  # 绝对路径
                    if file_path.endswith('.et'):
                        ket = wc.Dispatch("Ket.Application")
                        et = ket.Workbooks.Open(file_path)
                        # 将 .et 文件另存为 .xlsx 文件
                        xlsx_output_path = os.path.splitext(file_path)[0] + ".xlsx"
                        et.SaveAs(xlsx_output_path, 51)  # 使用参数 51 表示保存为 xlsx 格式
                        wb = openpyxl.load_workbook(xlsx_output_path)
                    elif file_path.endswith('.xlsx'):
                        wb = openpyxl.load_workbook(file_path)
                    if wb is not None:
                        for sheet_name in wb.sheetnames:
                            output.write("Sheet: " + sheet_name + '\n')
                            work_sheet = wb[sheet_name]

                            for ws_row in range(1, work_sheet.max_row + 1):
                                row_data = []
                                for ws_column in range(1, work_sheet.max_column + 1):
                                    cell_value = work_sheet.cell(row=ws_row, column=ws_column).value
                                    row_data.append(str(cell_value))

                                output_line = '\t'.join(row_data)  # 使用制表符分隔单元格数据，数据更整齐
                                output.write(output_line + '\n')
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                # # 对于.hiv文件，使用regipy库解析并输出注册表信息 （时间较长，可自行注释跳过，完成代码可行性测试）
                # elif file_path.endswith('.hiv'):
                #     file_path = os.path.abspath(file_path)  # 绝对路径
                #     hiv_output = parse_and_save_hiv(file_path)
                #     with open('output_hiv.txt', 'w', encoding='utf-8') as hivoutput:
                #         hivoutput.write('\n'.join(hiv_output) + '\n') # 单独写到一个txt文件里
                #     # output.write('\n'.join(hiv_output) + '\n') # 继续写到output里
                #     log.write("\nPreprocess " + file + " done \n") # 打印日志

                # 对于.wps格式文件，转换后缀成docx处理
                elif file_path.endswith('.wps'):
                    file_path.replace('.wps', '.docx')
                    file_path = os.path.abspath(file_path)  # 绝对路径
                    try:
                        text = ""
                        with open(file_path, 'r', encoding='utf-8') as f:
                            text += f.read()
                    except UnicodeDecodeError:
                        log.write(f"\nmaybe there are figures in {file} \n")
                    cleaned_text = clean_text(text)
                    output.write(cleaned_text + '\n')
                    log.write("\nPreprocess " + file + " done\n")  # 打印日志
                    pass

                # 对于.xml、.yml、.properties等文件，直接读取文件内容并写入输出文件
                elif file_path.endswith(('.pub', '.xml', '.yml', '.properties', '.py', '.sh', '.md', '.toml', '.rs')):
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        file_content = f.read()
                        output.write("File Path: " + file_path + '\n')
                        output.write("File Content:\n")
                        output.write(file_content + '\n')
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                # 处理图片文件, pytesseract用这个库
                elif file_path.endswith('.png') or file_path.endswith('.jpg'):
                    extracted_text = extract_all_text_from_image(file_path)
                    if extracted_text:
                        output.write("\nExtracted text from " + file + ": \n" + extracted_text + '\n')
                        log.write("Extraction from " + file + " done\n")
                    else:
                        log.write("No text found in photograph：" + file + '\n')

                # 处理.eml邮件文件
                elif file_path.endswith('.eml'):
                    with open(file_path, 'rb') as eml_file:
                        msg = email.message_from_bytes(eml_file.read())
                        # 解析邮件的其他信息，例如主题、发件人、收件人等
                        subject = msg.get('subject', 'No Subject')
                        from_address = msg.get('from', 'Unknown Sender')
                        to_addresses = msg.get('to', 'Unknown Recipient')
                        output.write("\nExtracted text from " + file_path + ": \n")
                        output.write("Email Subject: " + subject + '\n')
                        output.write("From: " + from_address + '\n')
                        output.write("To: " + to_addresses + '\n')
                        log.write(f"\nStart preprocessing .eml file: {file_path}\n")

                        for part in msg.walk():
                            content_type = part.get_content_type()
                            # 检查内容类型是否以 application/ 或 image/ 开头，表示这是一个附件。
                            if content_type.startswith('application/') or content_type.startswith('image/'):
                                filename = part.get_filename()
                                if filename:
                                    decoded_filename = decode_mime_word(filename)
                                    log.write(f"\n{file_path} has attachment: {decoded_filename}\n")
                                    attachment_data = part.get_payload(decode=True)
                                    attachment_path = os.path.join(root, decoded_filename)
                                    attachment_path = os.path.normpath(attachment_path)  # 规范化路径
                                    # print(attachment_path)
                                    with open(attachment_path, 'wb') as attachment_file:
                                        attachment_file.write(attachment_data)
                                    log.write("Attachment saved: " + attachment_path + '\n')
                                    # 根据附件后缀名调用相应的解析函数
                                    if decoded_filename.endswith('.txt'):
                                        with open(attachment_path, 'r', encoding='utf-8') as f:
                                            text = f.read()
                                        cleaned_text = clean_text(text)
                                        output.write(cleaned_text + '\n')
                                        log.write("Preprocess attachment " + decoded_filename + " done\n")
                                    elif decoded_filename.endswith('.xlsx'):
                                        wb = openpyxl.load_workbook(attachment_path)
                                        if wb is not None:
                                            for sheet_name in wb.sheetnames:
                                                output.write("Sheet: " + sheet_name + '\n')
                                                work_sheet = wb[sheet_name]
                                                for ws_row in range(1, work_sheet.max_row + 1):
                                                    row_data = []
                                                    for ws_column in range(1, work_sheet.max_column + 1):
                                                        cell_value = work_sheet.cell(row=ws_row, column=ws_column).value
                                                        row_data.append(str(cell_value))

                                                    output_line = '\t'.join(row_data)  # 使用制表符分隔单元格数据，数据更整齐
                                                    output.write(output_line + '\n')
                                        log.write("Preprocess attachment " + decoded_filename + " done\n")

                            elif content_type.startswith('text/'):
                                # 提取邮件正文的文本内容
                                email_text = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                                soup = BeautifulSoup(email_text, 'html.parser')
                                plaintext = soup.get_text()
                                cleaned_email_text = clean_text(plaintext)
                                output.write(cleaned_email_text + '\n')
                                log.write("Preprocess Email Text done\n")
                        log.write("Preprocess " + file + " done\n")

                elif file_path.endswith('.zip'):
                    pass

                # 最后处理没有后缀的文件
                elif file_path.endswith(''):
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        file_content = f.read()
                        output.write("File Path: " + file_path + '\n')
                        output.write("File Content:\n")
                        output.write(file_content + '\n')
                    log.write("\nPreprocess " + file + " done \n")  # 打印日志

                else:
                    log.write(f"\nUnsupported file format: {file_path} \n")


# 函数独立性测试，测试后自行注释掉
if __name__ == "__main__":
    # 假设解压后的数据在名为data的文件夹中
    data_folder = "data"
    output_file = "preprocessed_data.txt"
    preprocess_log = "preprocess_log.txt"
    preprocess_files_and_save(data_folder, output_file, preprocess_log)
    print("Preprocess all worked.")
